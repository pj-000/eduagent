import os
import asyncio
import concurrent.futures
import json
from typing import Dict, Any, List, Optional
from uuid import uuid4
from datetime import datetime
from pathlib import Path, PureWindowsPath, PurePosixPath
import logging
import traceback

try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None


RAG_EMBEDDING_MODEL = "BAAI/bge-small-zh-v1.5"
VECTOR_DIMENSION = 512  

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("faiss_toolkit")

try:
    from evoagentx.tools import Tool, Toolkit
    from evoagentx.tools.tool import ToolMetadata
except ImportError:
    try:
        from evoagentx.tools.tool import Tool, Toolkit, ToolMetadata
    except ImportError:
        class Tool:
            def __init__(self, name: str, description: str, inputs: Dict[str, Any], required: List[str]):
                if not name:
                    raise ValueError("Attribute name is required")
                self.name = name
                self.description = description or ""
                self.inputs = inputs or {}
                self.required = required or []
                self.metadata = ToolMetadata(name=name, description=description) if 'ToolMetadata' in locals() else None
            
            def __call__(self, *args, **kwargs):
                raise NotImplementedError("Tool __call__ method must be implemented")

        class Toolkit:
            def __init__(self, name: str, tools: List[Tool]):
                self.name = name
                self.tools = tools or []
                self.tool_map = {tool.name: tool for tool in self.tools}
            
            def get_tool(self, name: str) -> Optional[Tool]:
                return self.tool_map.get(name)

# 基础模块兼容
try:
    from evoagentx.core.module import BaseModule
except ImportError:
    class BaseModule:
        def __init__(self, **kwargs):
            pass

# 核心RAG相关导入
from evoagentx.rag.rag import RAGEngine
from evoagentx.rag.rag_config import RAGConfig
from evoagentx.rag.schema import Query, Document, Corpus, DocumentMetadata
from evoagentx.storages.base import StorageHandler
from evoagentx.storages.storages_config import StoreConfig

# 存储处理器兼容
try:
    from evoagentx.tools.storage_file import LocalStorageHandler
except ImportError:
    try:
        from evoagentx.storages import LocalStorageHandler
    except ImportError:
        class LocalStorageHandler:
            def __init__(self, base_path):
                self.base_path = base_path
                Path(base_path).mkdir(parents=True, exist_ok=True)
            
            def read(self, file_path):
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        return {"success": True, "content": f.read()}
                except Exception as e:
                    return {"success": False, "error": str(e)}

try:
    from evoagentx.tools.storage_handler import FileStorageHandler
except ImportError:
    class FileStorageHandler:
        pass

def _ensure_database_path(db_path: str) -> str:
    if not db_path:
        raise ValueError("Database path cannot be empty")
    
    path = Path(db_path).resolve()
    if path.exists() and path.is_dir():
        raise ValueError(f"Database path points to a directory: {db_path}")
    
    try:
        path.parent.mkdir(parents=True, exist_ok=True)
    except Exception as e:
        raise ValueError(f"Cannot create directory for database path {db_path}: {e}")
    
    if path.exists():
        logger.info(f"Found existing database at: {db_path}")
        try:
            import sqlite3
            conn = sqlite3.connect(str(path))
            conn.execute("SELECT name FROM sqlite_master WHERE type='table';")
            conn.close()
        except Exception as e:
            logger.warning(f"Database validation failed: {e}. Will create new database.")
            try:
                path.unlink()
            except Exception:
                pass
    else:
        logger.info(f"Database not found at: {db_path}. Will create new database.")
    
    return str(path)

def _create_default_storage_config(db_path: Optional[str] = None) -> StoreConfig:
    from evoagentx.storages.storages_config import StoreConfig, DBConfig, VectorStoreConfig
    
    if db_path is None:
        db_path = "./faiss_db.sqlite"
    
    validated_db_path = _ensure_database_path(db_path)
    index_cache_path = str(Path(validated_db_path).parent.resolve() / "index_cache")
    
    storage_config = StoreConfig(
        dbConfig=DBConfig(db_name="sqlite", path=validated_db_path),
        vectorConfig=VectorStoreConfig(
            vector_name="faiss",
            dimensions=VECTOR_DIMENSION,
            index_type="flat_l2"
        ),
        path=index_cache_path
    )
    
    Path(index_cache_path).mkdir(parents=True, exist_ok=True)
    return storage_config

def _create_default_rag_config() -> RAGConfig:
    from evoagentx.rag.rag_config import RAGConfig, EmbeddingConfig, ChunkerConfig
    
    return RAGConfig(
        embedding=EmbeddingConfig(
            provider="huggingface",
            model_name=RAG_EMBEDDING_MODEL 
        ),
        chunker=ChunkerConfig(
            chunk_size=500,
            chunk_overlap=50
        )
    )

def _make_json_serializable(obj: Any) -> Any:
    """递归将对象转换为JSON可序列化的格式"""
    if obj is None:
        return None
    elif isinstance(obj, (str, int, float, bool)):
        return obj
    elif isinstance(obj, datetime):
        return obj.isoformat()
    elif isinstance(obj, (Path, PureWindowsPath, PurePosixPath)):
        return str(obj)
    elif isinstance(obj, dict):
        return {k: _make_json_serializable(v) for k, v in obj.items()}
    elif isinstance(obj, (list, tuple, set)):
        return [_make_json_serializable(item) for item in obj]
    else:
        try:
            return str(obj)
        except:
            return None

def _clean_metadata_for_document(metadata: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
    """清理metadata，确保所有值都是JSON可序列化的"""
    if metadata is None:
        metadata = {}
    
    cleaned_metadata = {}
    for key, value in metadata.items():
        cleaned_metadata[key] = _make_json_serializable(value)
    
    if "insertion_time" not in cleaned_metadata:
        cleaned_metadata["insertion_time"] = datetime.now().isoformat()
    elif not isinstance(cleaned_metadata["insertion_time"], str):
        cleaned_metadata["insertion_time"] = _make_json_serializable(cleaned_metadata["insertion_time"])
    
    return cleaned_metadata

class FaissDatabase(BaseModule):
    def __init__(
        self,
        storage_config: StoreConfig,
        rag_config: RAGConfig,
        default_corpus_id: str = "default",
        default_index_type: str = "vector",
        storage_handler: StorageHandler = None,
        file_handler: FileStorageHandler = None,
        **kwargs
    ):
        super().__init__(**kwargs)
        self.storage_handler = StorageHandler(storageConfig=storage_config)
        
        try:
            from sentence_transformers import SentenceTransformer
            import numpy as np
            import evoagentx.rag.embeddings.huggingface_embedding as hf_emb

            if hasattr(hf_emb, 'SUPPORTED_MODELS'):
                if 'huggingface' in hf_emb.SUPPORTED_MODELS:
                    supported_list = hf_emb.SUPPORTED_MODELS['huggingface']
                    if RAG_EMBEDDING_MODEL not in supported_list:
                        supported_list.append(RAG_EMBEDDING_MODEL)
                        logger.info(f"🔓 HACK: Force-added '{RAG_EMBEDDING_MODEL}' to evoagentx whitelist.")
            
            logger.info(f"Loading local embedding model: {RAG_EMBEDDING_MODEL}...")
            local_embedding_model = SentenceTransformer(RAG_EMBEDDING_MODEL)
            local_embedding_model.max_seq_length = 512
            logger.info(f"Loaded local embedding model: {RAG_EMBEDDING_MODEL}")

            def custom_get_text_embeddings(texts: List[str]) -> List[List[float]]:
                embeddings = []
                for text in texts:
                    if not text or not text.strip():
                        embeddings.append(np.zeros(VECTOR_DIMENSION).tolist())
                    else:
                        emb = local_embedding_model.encode(
                            text,
                            normalize_embeddings=True,
                            show_progress_bar=False
                        )
                        embeddings.append(emb.tolist())
                return embeddings


            OriginalHuggingFaceEmbedding = hf_emb.HuggingFaceEmbedding

            class LocalBgeEmbedding(OriginalHuggingFaceEmbedding):
                """一个真正工作的HuggingFaceEmbedding子类，内部使用本地BGE模型。"""
                def __init__(self, model_name: str = RAG_EMBEDDING_MODEL, **kwargs):
                    super().__init__(model_name=RAG_EMBEDDING_MODEL, **kwargs)
                    
                    self._local_model = local_embedding_model
                    self.model_name = model_name
                    logger.info(f"🔧 LocalBgeEmbedding initialized with correct model: {model_name}")
                    
                    # 替换父类的内部方法
                    self._get_text_embeddings = custom_get_text_embeddings
                    self._get_text_embedding = lambda text: custom_get_text_embeddings([text])[0]

                def get_text_embedding_batch(self, texts, **kwargs):
                    return custom_get_text_embeddings(texts)
                
                def get_query_embedding(self, text, **kwargs):
                    return custom_get_text_embeddings([text])[0]

            hf_emb.HuggingFaceEmbedding = LocalBgeEmbedding
            logger.info("✅ Successfully patched HuggingFaceEmbedding class with LocalBgeEmbedding")

            self.rag_engine = RAGEngine(config=rag_config, storage_handler=self.storage_handler)
            logger.info("✅ RAGEngine initialized with local embedding bypass")
            
            try:
                logger.info(f"🔄 Attempting to load existing index for corpus: {default_corpus_id}...")
                self.rag_engine.load(corpus_id=default_corpus_id)
                if default_corpus_id in self.rag_engine.indices:
                    logger.info(f"✅ Successfully loaded index for corpus: {default_corpus_id}")
                else:
                    logger.info(f"ℹ️ Corpus {default_corpus_id} not found in storage (normal for fresh start)")
            except Exception as e:
                logger.warning(f"⚠️ Could not load existing index: {e}")
        
        except Exception as e:
            logger.error(f"❌ Failed to load local embedding model: {e}")
            logger.error(traceback.format_exc())
            raise
        
        if storage_handler is None:
            storage_handler = LocalStorageHandler(base_path=".")
        
        self.file_storage_handler = storage_handler
        self.default_corpus_id = default_corpus_id
        self.default_index_type = default_index_type
        logger.info(f"✅ Initialized FAISS database with corpus_id: {default_corpus_id}")
    
    def query(self, query: str, corpus_id: Optional[str] = None, top_k: int = 5, similarity_threshold: float = 0.0, metadata_filters: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        try:
            try:
                asyncio.get_running_loop()
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(self._query_sync, query, corpus_id, top_k, similarity_threshold, metadata_filters)
                    return future.result()
            except RuntimeError:
                return self._query_sync(query, corpus_id, top_k, similarity_threshold, metadata_filters)
        except Exception as e:
            logger.error(f"Query failed: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def _query_sync(self, query: str, corpus_id: Optional[str] = None, top_k: int = 5, similarity_threshold: float = 0.0, metadata_filters: Optional[Dict[str, Any]] = None) -> Dict[str, Any]:
        try:
            corpus_id = corpus_id or self.default_corpus_id
            
            if corpus_id not in self.rag_engine.indices:
                return {
                    "success": True,
                    "data": {
                        "query": query,
                        "corpus_id": corpus_id,
                        "total_results": 0,
                        "results": []
                    }
                }
            
            query_obj = Query(
                query_str=query,
                top_k=top_k,
                similarity_cutoff=similarity_threshold,
                metadata_filters=metadata_filters
            )
            
            results = self.rag_engine.query(query_obj, corpus_id=corpus_id)
            
            formatted_results = {
                "query": query,
                "corpus_id": corpus_id,
                "total_results": 0,
                "results": []
            }
            
            if results and results.corpus and results.corpus.chunks:
                chunks = results.corpus.chunks
                formatted_results["total_results"] = len(chunks)
                
                for i, chunk in enumerate(chunks):
                    score = results.scores[i] if (results.scores and i < len(results.scores)) else 0.0
                    formatted_results["results"].append({
                        "chunk_id": chunk.chunk_id,
                        "content": chunk.text,
                        "score": score,
                        "metadata": chunk.metadata.model_dump() if chunk.metadata else {},
                        "doc_id": chunk.metadata.doc_id if (chunk.metadata and hasattr(chunk.metadata, 'doc_id')) else None
                    })
            
            return {"success": True, "data": formatted_results}
        
        except Exception as e:
            logger.error(f"Sync query failed: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def _is_file_path(self, text: str) -> bool:
        # 扩展支持的文件类型
        path_indicators = ['/', '\\', '.txt', '.pdf', '.md', '.doc', '.docx', '.ppt', '.pptx', '.csv', '.json', '.xml']
        return any(indicator in text for indicator in path_indicators) and os.path.exists(text)
    
    def _process_file_path(self, file_path: str, doc_index: int, metadata: Optional[Dict[str, Any]] = None) -> List[Document]:
        try:
            try:
                asyncio.get_running_loop()
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(self._process_file_path_sync, file_path, doc_index, metadata)
                    return future.result()
            except RuntimeError:
                return self._process_file_path_sync(file_path, doc_index, metadata)
        except Exception as e:
            logger.error(f"Failed to process file {file_path}: {str(e)}")
            doc_metadata = _clean_metadata_for_document(metadata)
            doc_metadata.update({
                "doc_index": doc_index,
                "source_file": file_path,
                "error": str(e)
            })
            document_metadata = DocumentMetadata(**doc_metadata)
            return [Document(
                text=f"Error reading file {file_path}: {str(e)}",
                metadata=document_metadata,
                doc_id=str(uuid4())
            )]
    
    def _process_file_path_sync(self, file_path: str, doc_index: int, metadata: Optional[Dict[str, Any]] = None) -> List[Document]:
        try:
            file_content = ""
            file_path_obj = Path(file_path).resolve()
            ext = file_path_obj.suffix.lower()
            
            if ext == '.pdf':
                if pdfplumber is None:
                    raise ImportError("pdfplumber is not installed. Run: pip install pdfplumber")
                try:
                    text_parts = []
                    with pdfplumber.open(file_path_obj) as pdf:
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                paragraphs = page_text.split('\n\n')
                                # 清洗空白字符并重组
                                text_parts.append("\n".join(p.strip() for p in paragraphs if p.strip()))
                    file_content = "\n\n".join(text_parts)
                except Exception as e:
                    logger.error(f"Error reading PDF {file_path}: {e}")
                    raise
            
            elif ext == '.docx':
                if DocxDocument is None:
                    raise ImportError("python-docx is not installed. Run: pip install python-docx")
                try:
                    doc = DocxDocument(file_path_obj)
                    paragraphs = [para.text.strip() for para in doc.paragraphs if para.text.strip()]
                    file_content = "\n".join(paragraphs)
                except Exception as e:
                    logger.error(f"Error reading DOCX {file_path}: {e}")
                    raise

            elif ext in ['.ppt', '.pptx']:
                if Presentation is None:
                    raise ImportError("python-pptx is not installed. Run: pip install python-pptx")
                try:
                    prs = Presentation(file_path_obj)
                    text_runs = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.text.strip():
                                        text_runs.append(run.text.strip())
                    file_content = '\n'.join(text_runs)
                except Exception as e:
                    logger.error(f"Error reading PPT {file_path}: {e}")
                    raise

            else:
                try:
                    with open(file_path_obj, 'r', encoding='utf-8', errors='ignore') as f:
                        file_content = f.read().strip()
                except Exception as e:
                     logger.error(f"Error reading text file {file_path}: {e}")
                     raise
            
            # 如果内容为空，给一个提示
            if not file_content:
                logger.warning(f"File {file_path} extracted content is empty.")
                file_content = f"[Empty content extracted from {file_path_obj.name}]"

            # 构建 Document 对象
            doc_metadata = _clean_metadata_for_document(metadata)
            doc_metadata.update({
                "doc_index": doc_index,
                "source_file": str(file_path_obj),
                "file_type": ext,
                "file_size": len(file_content)
            })
            document_metadata = DocumentMetadata(**doc_metadata)
            
            return [Document(
                text=file_content,
                metadata=document_metadata,
                doc_id=str(uuid4())
            )]
        
        except Exception as e:
            logger.error(f"Failed to process file {file_path} in sync mode: {str(e)}")
            doc_metadata = _clean_metadata_for_document(metadata)
            doc_metadata.update({
                "doc_index": doc_index,
                "source_file": file_path,
                "error": str(e)
            })
            document_metadata = DocumentMetadata(**doc_metadata)
            return [Document(
                text=f"Error reading file {file_path}: {str(e)}",
                metadata=document_metadata,
                doc_id=str(uuid4())
            )]
    
    def insert(self, documents: list, corpus_id: Optional[str] = None, metadata: Optional[Dict[str, Any]] = None, batch_size: int = 100) -> Dict[str, Any]:
        try:
            corpus_id = corpus_id or self.default_corpus_id
            processed_docs = []
            file_paths_processed = []
            
            for i, doc in enumerate(documents):
                if isinstance(doc, str):
                    if self._is_file_path(doc):
                        file_docs = self._process_file_path(doc, i, metadata)
                        processed_docs.extend(file_docs)
                        file_paths_processed.append(doc)
                    else:
                        doc_metadata = _clean_metadata_for_document(metadata)
                        doc_metadata.update({
                            "doc_index": i
                        })
                        document_metadata = DocumentMetadata(**doc_metadata)
                        processed_docs.append(Document(
                            text=doc,
                            metadata=document_metadata,
                            doc_id=str(uuid4())
                        ))
                elif isinstance(doc, dict):
                    doc_metadata = _clean_metadata_for_document(metadata)
                    doc_metadata.update(_clean_metadata_for_document(doc.get("metadata", {})))
                    doc_metadata.update({
                        "doc_index": i
                    })
                    document_metadata = DocumentMetadata(**doc_metadata)
                    processed_docs.append(Document(
                        text=doc.get("text", ""),
                        metadata=document_metadata,
                        doc_id=doc.get("doc_id", str(uuid4()))
                    ))
            
            corpus = Corpus(corpus_id=corpus_id)
            total_processed = 0
            
            for i in range(0, len(processed_docs), batch_size):
                batch = processed_docs[i:i+batch_size]
                batch_corpus = self.rag_engine.chunker.chunk(batch)
                batch_corpus.corpus_id = corpus_id
                
                self.rag_engine.add(self.default_index_type, batch_corpus, corpus_id=corpus_id)
                corpus.chunks.extend(batch_corpus.chunks)
                
                total_processed += len(batch)
                logger.info(f"Processed batch {i//batch_size + 1}, total processed: {total_processed}")
            
            for chunk in corpus.chunks:
                if chunk.metadata:
                    try:
                        json.dumps(chunk.metadata.model_dump())
                    except Exception as e:
                        logger.warning(f"Chunk {chunk.chunk_id} metadata not JSON serializable, cleaning...")
                        cleaned_metadata = _clean_metadata_for_document(chunk.metadata.model_dump())
                        chunk.metadata = DocumentMetadata(**cleaned_metadata)
            
            self.rag_engine.save(corpus_id=corpus_id, index_type=self.default_index_type)
            
            result = {
                "corpus_id": corpus_id,
                "documents_inserted": len(documents),
                "chunks_created": len(corpus.chunks),
                "total_processed": total_processed,
                "file_paths_processed": file_paths_processed
            }
            
            return {"success": True, "data": result}
        
        except Exception as e:
            logger.error(f"Insert failed: {str(e)}")
            logger.error(traceback.format_exc())
            return {"success": False, "error": str(e)}
    
    def delete(self, corpus_id: Optional[str] = None, doc_ids: Optional[List[str]] = None, metadata_filters: Optional[Dict[str, Any]] = None, clear_all: bool = False) -> Dict[str, Any]:
        try:
            corpus_id = corpus_id or self.default_corpus_id
            
            if clear_all:
                self.rag_engine.clear(corpus_id=corpus_id)
                return {
                    "success": True,
                    "data": {"operation": "clear_all", "corpus_id": corpus_id}
                }
            
            if corpus_id not in self.rag_engine.indices:
                return {
                    "success": True,
                    "data": {
                        "operation": "selective_delete",
                        "corpus_id": corpus_id,
                        "message": "Corpus not found"
                    }
                }
            
            if doc_ids or metadata_filters:
                cleaned_filters = _clean_metadata_for_document(metadata_filters) if metadata_filters else None
                self.rag_engine.delete(
                    corpus_id=corpus_id,
                    index_type=self.default_index_type,
                    node_ids=doc_ids,
                    metadata_filters=cleaned_filters
                )
                return {
                    "success": True,
                    "data": {
                        "corpus_id": corpus_id,
                        "operation": "selective_delete",
                        "doc_ids": doc_ids,
                        "metadata_filters": cleaned_filters
                    }
                }
            else:
                return {
                    "success": True,
                    "data": {
                        "operation": "selective_delete",
                        "corpus_id": corpus_id,
                        "message": "No criteria provided"
                    }
                }
        
        except Exception as e:
            logger.error(f"Delete failed: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def list_corpora(self) -> Dict[str, Any]:
        try:
            corpora = []
            for corpus_id, indices in self.rag_engine.indices.items():
                corpus_info = {
                    "corpus_id": corpus_id,
                    "index_types": list(indices.keys()),
                    "retrievers": list(self.rag_engine.retrievers.get(corpus_id, {}).keys())
                }
                corpora.append(corpus_info)
            
            return {
                "success": True,
                "data": {"corpora": corpora, "total": len(corpora)}
            }
        
        except Exception as e:
            logger.error(f"List corpora failed: {str(e)}")
            return {"success": False, "error": str(e)}
    
    def get_stats(self, corpus_id: Optional[str] = None) -> Dict[str, Any]:
        try:
            if corpus_id:
                corpus_id = corpus_id or self.default_corpus_id
                stats = {
                    "corpus_id": corpus_id,
                    "exists": corpus_id in self.rag_engine.indices,
                    "index_types": list(self.rag_engine.indices.get(corpus_id, {}).keys()),
                    "retrievers": list(self.rag_engine.retrievers.get(corpus_id, {}).keys())
                }
                
                if corpus_id in self.rag_engine.indices:
                    vector_index = self.rag_engine.indices[corpus_id].get(self.default_index_type)
                    if vector_index and hasattr(vector_index, 'get_index'):
                        try:
                            index = vector_index.get_index()
                            if hasattr(index, 'vector_store') and hasattr(index.vector_store, 'faiss_index'):
                                stats["vector_count"] = index.vector_store.faiss_index.ntotal
                                stats["dimensions"] = index.vector_store.faiss_index.d
                        except Exception:
                            pass
                
                return {"success": True, "data": stats}
            else:
                stats = {
                    "total_corpora": len(self.rag_engine.indices),
                    "corpora": list(self.rag_engine.indices.keys()),
                    "embedding_model": RAG_EMBEDDING_MODEL,
                    "vector_store_type": self.rag_engine.storage_handler.storageConfig.vectorConfig.vector_name 
                    if (self.rag_engine.storage_handler and self.rag_engine.storage_handler.storageConfig and self.rag_engine.storage_handler.storageConfig.vectorConfig) 
                    else None
                }
                return {"success": True, "data": stats}
        
        except Exception as e:
            logger.error(f"Get stats failed: {str(e)}")
            return {"success": False, "error": str(e)}

class FaissQueryTool(Tool):
    def __init__(self, faiss_database: FaissDatabase = None):
        super().__init__(
            name="faiss_query",
            description="Query the FAISS vector database using semantic search",
            inputs={
                "query": {"type": "string", "description": "Search query text"},
                "corpus_id": {"type": "string", "description": "Optional corpus ID"},
                "top_k": {"type": "integer", "description": "Number of results", "default": 5},
                "similarity_threshold": {"type": "number", "description": "Minimum similarity", "default": 0.0},
                "metadata_filters": {"type": "object", "description": "Optional metadata filters"}
            },
            required=["query"]
        )
        self.faiss_database = faiss_database
    
    def __call__(self, query: str, corpus_id: str = None, top_k: int = 5, similarity_threshold: float = 0.0, metadata_filters: dict = None, **kwargs) -> Dict[str, Any]:
        return self.faiss_database.query(
            query=query,
            corpus_id=corpus_id,
            top_k=top_k,
            similarity_threshold=similarity_threshold,
            metadata_filters=metadata_filters
        )

class FaissInsertTool(Tool):
    def __init__(self, faiss_database: FaissDatabase = None):
        super().__init__(
            name="faiss_insert",
            description="Insert documents or files into the FAISS vector database.",
            inputs={
                "documents": {"type": "array", "description": "List of strings (text/paths) or objects"},
                "corpus_id": {"type": "string", "description": "Optional corpus ID"},
                "metadata": {"type": "object", "description": "Common metadata"},
                "batch_size": {"type": "integer", "description": "Batch size", "default": 100}
            },
            required=["documents"]
        )
        self.faiss_database = faiss_database
    
    def __call__(self, documents: list, corpus_id: str = None, metadata: dict = None, batch_size: int = 100, **kwargs) -> Dict[str, Any]:
        return self.faiss_database.insert(
            documents=documents,
            corpus_id=corpus_id,
            metadata=metadata,
            batch_size=batch_size
        )

class FaissDeleteTool(Tool):
    def __init__(self, faiss_database: FaissDatabase = None):
        super().__init__(
            name="faiss_delete",
            description="Delete documents from FAISS vector database",
            inputs={
                "corpus_id": {"type": "string", "description": "Optional corpus ID"},
                "doc_ids": {"type": "array", "description": "List of doc IDs to delete"},
                "metadata_filters": {"type": "object", "description": "Metadata filters"},
                "clear_all": {"type": "boolean", "description": "Clear all documents", "default": False}
            },
            required=[]
        )
        self.faiss_database = faiss_database
    
    def __call__(self, corpus_id: str = None, doc_ids: list = None, metadata_filters: dict = None, clear_all: bool = False, **kwargs) -> Dict[str, Any]:
        return self.faiss_database.delete(
            corpus_id=corpus_id,
            doc_ids=doc_ids,
            metadata_filters=metadata_filters,
            clear_all=clear_all
        )

class FaissListTool(Tool):
    def __init__(self, faiss_database: FaissDatabase = None):
        super().__init__(
            name="faiss_list",
            description="List all available corpora in FAISS database",
            inputs={},
            required=[]
        )
        self.faiss_database = faiss_database
    
    def __call__(self, **kwargs) -> Dict[str, Any]:
        return self.faiss_database.list_corpora()

class FaissStatsTool(Tool):
    def __init__(self, faiss_database: FaissDatabase = None):
        super().__init__(
            name="faiss_stats",
            description="Get database statistics and information",
            inputs={
                "corpus_id": {"type": "string", "description": "Optional corpus ID"}
            },
            required=[]
        )
        self.faiss_database = faiss_database
    
    def __call__(self, corpus_id: str = None, **kwargs) -> Dict[str, Any]:
        return self.faiss_database.get_stats(corpus_id=corpus_id)

class FaissToolkit(Toolkit):
    def __init__(self, name: str = "FaissToolkit", storage_config: Optional[StoreConfig] = None, rag_config: Optional[RAGConfig] = None, default_corpus_id: str = "default", default_index_type: str = "vector", db_path: Optional[str] = None, storage_handler: StorageHandler = None, file_handler: FileStorageHandler = None, **kwargs):
        if storage_config is None:
            storage_config = _create_default_storage_config(db_path)
        
        if rag_config is None:
            rag_config = _create_default_rag_config()
        
        faiss_database = FaissDatabase(
            storage_config=storage_config,
            rag_config=rag_config,
            default_corpus_id=default_corpus_id,
            default_index_type=default_index_type,
            storage_handler=storage_handler,
            file_handler=file_handler
        )
        
        tools = [
            FaissQueryTool(faiss_database),
            FaissInsertTool(faiss_database),
            FaissDeleteTool(faiss_database),
            FaissListTool(faiss_database),
            FaissStatsTool(faiss_database)
        ]
        
        super().__init__(name=name, tools=tools)
        self.faiss_database = faiss_database
    
    def get_database(self) -> FaissDatabase:
        return self.faiss_database
    
    def get_tool(self, name: str) -> Optional[Tool]:
        return super().get_tool(name)