"""
幻灯片创建器
负责创建幻灯片并添加格式化内容
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from content_extractor import split_content_by_length, extract_content_elements_enhanced, split_content_by_order
from content_formatter import apply_code_syntax_highlighting, distribute_code_to_placeholders, format_table_content, format_text_with_bold
from layout_selector import choose_layout_enhanced
from utils import apply_paragraph_format_respecting_template, apply_default_content_format
from config import *

prs = Presentation(TMPL_FILE)
layout_map = {ly.name: ly for ly in prs.slide_layouts}
# 确保 layout_selector 使用同一个 Presentation 实例，避免跨实例引用导致的损坏
import layout_selector as _layout_selector
_layout_selector.layout_map = layout_map
_layout_selector.prs = prs

def add_content_to_slide_enhanced(slide, elements, title_text="", subtitle_text="", is_toc=False, toc_items=None, is_title_page=False):
    """
    增强版幻灯片内容添加，支持代码高亮、表格格式化和字体控制
    重点修复：确保普通文本不被当作列表处理
    新增：支持标题页处理
    """
    # 确保 elements 包含所有必要的键，防止 KeyError
    elements.setdefault('plain_text', [])
    elements.setdefault('bullets', [])
    elements.setdefault('numbered_lists', [])
    elements.setdefault('tables', [])
    elements.setdefault('code_blocks', [])
    elements.setdefault('quotes', [])
    
    # 获取所有文本占位符，并按位置排序
    text_placeholders = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and hasattr(shape, 'placeholder_format'):
            text_placeholders.append(shape)
    
    # 按占位符的索引ID排序，确保标题占位符在前
    try:
        text_placeholders.sort(key=lambda x: x.placeholder_format.idx)
    except:
        # 如果无法按索引排序，按垂直位置排序（上面的是标题）
        text_placeholders.sort(key=lambda x: x.top)
    
    layout_name = slide.slide_layout.name
    print(f"使用版式: {layout_name}")
    print(f"找到 {len(text_placeholders)} 个文本占位符")
    
    # 调试信息：显示占位符信息
    for i, placeholder in enumerate(text_placeholders):
        try:
            idx = placeholder.placeholder_format.idx
            print(f"  占位符{i}: idx={idx}, top={placeholder.top}, left={placeholder.left}")
        except:
            print(f"  占位符{i}: top={placeholder.top}, left={placeholder.left}")
    
    # 设置标题 - 使用第一个文本占位符
    title_placeholder = None
    if text_placeholders and title_text:
        title_placeholder = text_placeholders[0]
        
        # 如果有子标题，组合显示
        if subtitle_text:
            title_placeholder.text = f"{title_text}\n{subtitle_text}"
        else:
            title_placeholder.text = title_text
            
        print(f"  设置标题到占位符0: {title_text}")
            
        # 设置标题样式 - 特殊处理目录页标题
        if hasattr(title_placeholder, 'text_frame'):
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    # 如果是目录页，使用专门的字体大小配置
                    if is_toc:
                        run.font.size = Pt(TOC_TITLE_FONT_SIZE)
                        print(f"  设置目录标题字体大小: {TOC_TITLE_FONT_SIZE}pt")
                    else:
                        run.font.size = Pt(24)
    
    # 获取内容占位符 - 跳过第一个(标题)占位符
    content_placeholders = text_placeholders[1:] if len(text_placeholders) > 1 else []
    print(f"  内容占位符数量: {len(content_placeholders)}")
    
    # 处理目录页
    if is_toc and toc_items and content_placeholders:
        text_frame = content_placeholders[0].text_frame
        
        # 保存原始段落格式（模板中的设置）
        original_paragraph = None
        if text_frame.paragraphs:
            original_paragraph = text_frame.paragraphs[0]
        
        text_frame.clear()
        
        print(f"  创建目录页，目录项数量: {len(toc_items)}")
        
        for i, toc_item in enumerate(toc_items):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # 添加目录项，不再自动添加数字编号
            p.text = toc_item  # 直接使用原始目录项，保留"一、"、"二、"等编号
            p.alignment = PP_ALIGN.LEFT
            
            # 设置目录项字体样式 - 优先使用模板设置
            for run in p.runs:
                # 尝试从原始段落获取字体设置
                if original_paragraph and original_paragraph.runs:
                    if original_paragraph.runs[0].font.size:
                        run.font.size = original_paragraph.runs[0].font.size
                    else:
                        run.font.size = Pt(TOC_ITEMS_FONT_SIZE)  # 使用配置的目录项字体大小
                    
                    if original_paragraph.runs[0].font.name:
                        run.font.name = original_paragraph.runs[0].font.name
                    else:
                        run.font.name = '微软雅黑'  # 默认值
                else:
                    # 如果没有原始段落信息，使用配置的目录项字体大小
                    run.font.size = Pt(TOC_ITEMS_FONT_SIZE)
                    run.font.name = '微软雅黑'
            
            # 使用统一的段落格式处理函数，确保遵循模板设置
            apply_paragraph_format_respecting_template(p, original_paragraph)
            
            # 如果模板没有设置段落格式，应用默认的目录格式
            if original_paragraph is None:
                apply_default_content_format(p, "normal")
        
        return
    
    # 处理标题页
    if is_title_page:
        print(f"  创建标题页，标题: {title_text}")
        # 标题页通常只有主标题和副标题，不需要额外的内容处理
        # 标题已经在上面的代码中设置了，这里可以添加特殊的标题页样式
        if title_placeholder and hasattr(title_placeholder, 'text_frame'):
            for paragraph in title_placeholder.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(36)  # 标题页使用更大的字体
                    run.font.name = '微软雅黑'
        
        # 如果有副标题（课程标题页）
        if subtitle_text and content_placeholders:
            content_frame = content_placeholders[0].text_frame
            content_frame.clear()
            p = content_frame.paragraphs[0]
            p.text = subtitle_text
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.size = Pt(20)
                run.font.name = '微软雅黑'
                run.font.italic = True
        
        return
    
    # 根据不同版式处理内容
    if layout_name == 'TitleOnly':
        return
    elif layout_name in ['CodeSlide', 'CodeSlide1', 'CodeSlide2'] and elements['code_blocks']:
        if content_placeholders:
            # 修复：处理所有代码块，不只是第一个
            all_code_content = []
            for i, code_block in enumerate(elements['code_blocks']):
                if i > 0:
                    all_code_content.append("")  # 在代码块之间添加空行分隔
                    all_code_content.append(f"# --- 代码块 {i+1} ---")  # 添加分隔注释
                    all_code_content.append("")
                all_code_content.extend(code_block.split('\n'))
            
            combined_code_content = '\n'.join(all_code_content)
            print(f"  将 {len(elements['code_blocks'])} 个代码块合并分布到 {len(content_placeholders)} 个占位符中")
            # 使用新的代码分布函数
            distribute_code_to_placeholders(combined_code_content, content_placeholders, layout_name)
        else:
            print("  警告: 代码版式缺少内容占位符")
        return
    
    elif layout_name == 'TableSlide' and elements['tables']:
        if content_placeholders:
            table_content = elements['tables'][0]
            print(f"  将表格内容放入内容占位符")
            format_table_content(table_content, content_placeholders[0].text_frame)
        else:
            print("  警告: TableSlide版式缺少内容占位符")
        return
    
    elif layout_name == 'QuoteSlide' and elements['quotes']:
        if content_placeholders:
            quote_content = elements['quotes'][0]
            text_frame = content_placeholders[0].text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = f'"{quote_content}"'
            for run in p.runs:
                run.font.italic = True
        return
    
    elif layout_name == 'ImageLeftTextRight':
        if content_placeholders:
            text_frame = content_placeholders[-1].text_frame
            text_frame.clear()
            
            content_added = False
            
            # 首先添加普通文本（不加列表符号）
            for text in elements['plain_text']:
                if text.strip() and len(text.strip()) > 3:
                    if content_added:
                        p = text_frame.add_paragraph()
                    else:
                        p = text_frame.paragraphs[0]
                        content_added = True
                    
                    # 使用新的粗体处理函数
                    format_text_with_bold(p, text)
            
            # 然后添加列表项（使用粗体处理函数）
            for bullet in elements['bullets']:
                if content_added:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                    content_added = True
                
                # 使用新的粗体处理函数
                format_text_with_bold(p, bullet)
            
            for i, item in enumerate(elements['numbered_lists'], 1):
                if content_added:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                    content_added = True
                
                # 使用新的粗体处理函数
                format_text_with_bold(p, f"{i}. {item}")
        return
    
    elif layout_name in ['OneBullet', 'OneBullet1', 'OneBullet2', 'OneBullet3',
                         'TwoBullets', 'TwoBullets1', 'TwoBullets2', 'TwoBullets3', 
                         'ThreeBullets', 'ThreeBullets1', 'ThreeBullets2', 'ThreeBullets3',
                         'FourBullets', 'FourBullets1', 'FourBullets2', 'FourBullets3']:
        # 根据版式名称确定最大要点数
        if layout_name.startswith('OneBullet'):
            max_bullets = 1
        elif layout_name.startswith('TwoBullets'):
            max_bullets = 2
        elif layout_name.startswith('ThreeBullets'):
            max_bullets = 3
        elif layout_name.startswith('FourBullets'):
            max_bullets = 4
        else:
            max_bullets = 4  # 默认值
        
        # 合并所有内容作为要点
        all_bullets = []
        
        # 添加普通文本（作为要点，但不加bullet符号）
        for text in elements['plain_text']:
            if text.strip() and len(text.strip()) > 3:
                all_bullets.append(text)
        
        # 添加真正的bullet要点（不加符号）
        for bullet in elements['bullets']:
            all_bullets.append(bullet)
            
        # 添加数字列表
        for i, item in enumerate(elements['numbered_lists'], 1):
            all_bullets.append(f"{i}. {item}")
        
        for i, bullet in enumerate(all_bullets[:max_bullets]):
            if i < len(content_placeholders):
                text_frame = content_placeholders[i].text_frame
                # 保留模板原始段落格式信息
                original_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
                text_frame.clear()
                p = text_frame.paragraphs[0]
                
                # 使用新的粗体处理函数
                format_text_with_bold(p, bullet)
                
                # 遵循模板段落格式
                apply_paragraph_format_respecting_template(p, original_paragraph)
                if original_paragraph is None:
                    apply_default_content_format(p, "bullet")
        return
    
    # 默认处理 - 修复关键部分，避免所有内容都变成列表
    if content_placeholders:
        text_frame = content_placeholders[0].text_frame
        
        # 保存原始段落格式
        original_paragraph = None
        if text_frame.paragraphs:
            original_paragraph = text_frame.paragraphs[0]
        
        text_frame.clear()
        
        # 处理不同类型的内容，保持原始格式
        content_added = False
        
        # 首先添加普通文本段落（不加列表符号）
        for text in elements['plain_text']:
            if text.strip() and len(text.strip()) > 3:
                if content_added:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                    content_added = True
                
                # 使用新的粗体处理函数
                format_text_with_bold(p, text)
                # 使用模板格式或应用默认格式
                apply_paragraph_format_respecting_template(p, original_paragraph)
                if original_paragraph is None:
                    apply_default_content_format(p, "normal")
        
        # 然后添加要点列表（不带bullet符号）
        if elements['bullets']:
            # 添加一个空行分隔
            if content_added:
                p = text_frame.add_paragraph()
                p.text = ""
                if original_paragraph and hasattr(original_paragraph, 'space_after'):
                    p.space_after = Pt(4)
                content_added = True
            
            for bullet in elements['bullets']:
                if content_added:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                    content_added = True
                
                # 使用新的粗体处理函数
                format_text_with_bold(p, bullet)
                
                # 使用模板格式或应用默认格式
                apply_paragraph_format_respecting_template(p, original_paragraph)
                if original_paragraph is None:
                    apply_default_content_format(p, "bullet")
        
        # 添加数字列表
        if elements['numbered_lists']:
            # 添加一个空行分隔
            if content_added and elements['bullets']:
                p = text_frame.add_paragraph()
                p.text = ""
                if original_paragraph and hasattr(original_paragraph, 'space_after'):
                    p.space_after = Pt(4)
                content_added = True
            
            for i, item in enumerate(elements['numbered_lists'], 1):
                if content_added:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                    content_added = True
                
                # 使用新的粗体处理函数
                format_text_with_bold(p, f"{i}. {item}")
                
                # 使用模板格式或应用默认格式
                apply_paragraph_format_respecting_template(p, original_paragraph)
                if original_paragraph is None:
                    apply_default_content_format(p, "numbered")
        
        # 添加引用
        if elements['quotes']:
            # 添加一个空行分隔
            if content_added:
                p = text_frame.add_paragraph()
                p.text = ""
                if original_paragraph and hasattr(original_paragraph, 'space_after'):
                    p.space_after = Pt(4)
                content_added = True
            
            for quote in elements['quotes']:
                if content_added:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                    content_added = True
                
                p.text = f'"{quote}"'
                for run in p.runs:
                    run.font.italic = True
                
                # 使用模板格式或应用默认格式
                apply_paragraph_format_respecting_template(p, original_paragraph)
                if original_paragraph is None:
                    apply_default_content_format(p, "quote")
        
        # 添加代码块（当使用混合版式时）
        if elements['code_blocks']:
            # 添加一个空行分隔
            if content_added:
                p = text_frame.add_paragraph()
                p.text = ""
                if original_paragraph and hasattr(original_paragraph, 'space_after'):
                    p.space_after = Pt(6)
                content_added = True
            
            for code_block in elements['code_blocks']:
                # 添加代码块标题（可选）
                if content_added:
                    p = text_frame.add_paragraph()
                else:
                    p = text_frame.paragraphs[0]
                    content_added = True
                
                p.text = "代码示例："
                for run in p.runs:
                    run.font.bold = True
                    run.font.size = Pt(14)
                
                # 添加代码内容
                code_lines = code_block.split('\n')
                for line_num, code_line in enumerate(code_lines):
                    p = text_frame.add_paragraph()
                    apply_code_syntax_highlighting(p, code_line)
                    apply_paragraph_format_respecting_template(p, original_paragraph)
                
                content_added = True

def create_slides_optimized(slides_data):
    """
    优化版幻灯片创建流程，修复了混合内容的顺序问题
    """

    def get_chunk_title(base_title, chunk, chunk_idx, total_chunks):
        """根据内容类型决定幻灯片标题，避免误加“代码示例”"""
        has_non_code_content = any([
            len(chunk['bullets']) > 0,
            len(chunk['numbered_lists']) > 0,
            len(chunk['plain_text']) > 0,
            len(chunk['tables']) > 0,
            len(chunk['quotes']) > 0
        ])

        if total_chunks > 1:
            return f"{base_title} - 第{chunk_idx+1}部分"

        if chunk['code_blocks'] and not has_non_code_content:
            return f"{base_title} - 代码示例"

        return base_title

    total_slides_created = 0
    for slide_data in slides_data:
        main_title = slide_data['main_title']
        content = slide_data['content']
        sections = slide_data['sections']
        is_toc = slide_data.get('is_toc', False)
        is_title_page = slide_data.get('is_title_page', False)
        toc_items = slide_data.get('toc_items', [])
        
        print(f"\n处理章节: {main_title}")
        
        # 处理目录页
        if is_toc:
            layout = choose_layout_enhanced({}, is_toc=True)
            slide = prs.slides.add_slide(layout)
            add_content_to_slide_enhanced(slide, {}, main_title, "", is_toc=True, toc_items=toc_items)
            total_slides_created += 1
            print(f"  创建目录页: {len(toc_items)} 个目录项")
            continue
        
        # 处理标题页
        if is_title_page:
            layout = choose_layout_enhanced({}, is_title_page=True)
            slide = prs.slides.add_slide(layout)
            subtitle = ""  # 标题页不需要副标题
            add_content_to_slide_enhanced(slide, {}, main_title, subtitle, is_title_page=True)
            total_slides_created += 1
            print(f"  创建标题页: {main_title}")
            continue
        
        # 1. 为一级标题创建标题页（可选）
        if CREATE_TITLE_SLIDES:
            if content:
                # 有内容的一级标题，提取内容元素
                elements = extract_content_elements_enhanced(content)
                
                # 检查是否存在混合内容（要点+代码块）
                has_bullets_or_text = len(elements['bullets']) > 0 or len(elements['numbered_lists']) > 0 or len(elements['plain_text']) > 0
                has_code_blocks = len(elements['code_blocks']) > 0
                
                if has_bullets_or_text and has_code_blocks:
                    print(f"  检测到一级标题混合内容，将按原始顺序创建幻灯片")
                    
                    # 使用新的按顺序分割函数
                    ordered_chunks = split_content_by_order(elements)
                    
                    for chunk_idx, chunk in enumerate(ordered_chunks):
                        layout = choose_layout_enhanced(chunk)
                        slide = prs.slides.add_slide(layout)
                        
                        chunk_title = get_chunk_title(main_title, chunk, chunk_idx, len(ordered_chunks))
                        
                        add_content_to_slide_enhanced(slide, chunk, chunk_title, "")
                        total_slides_created += 1
                        
                        content_desc = []
                        if chunk['bullets']:
                            content_desc.append(f"{len(chunk['bullets'])} 要点")
                        if chunk['numbered_lists']:
                            content_desc.append(f"{len(chunk['numbered_lists'])} 数字列表")
                        if chunk['plain_text']:
                            content_desc.append(f"{len(chunk['plain_text'])} 文本段落")
                        if chunk['code_blocks']:
                            content_desc.append(f"{len(chunk['code_blocks'])} 代码块")
                        
                        print(f"    创建一级标题混合内容页 {chunk_idx+1}/{len(ordered_chunks)}: {', '.join(content_desc)}")
                else:
                    # 非混合内容，使用按顺序分割以保持原始顺序
                    chunks = split_content_by_order(elements)
                    
                    for chunk_idx, chunk in enumerate(chunks):
                        layout = choose_layout_enhanced(chunk)
                        slide = prs.slides.add_slide(layout)
                        
                        subtitle = f"第{chunk_idx+1}部分" if len(chunks) > 1 else ""
                        add_content_to_slide_enhanced(slide, chunk, main_title, subtitle)
                        total_slides_created += 1
                        print(f"  创建一级标题内容页 {chunk_idx+1}/{len(chunks)}: {len(chunk['plain_text'])} 文本段落, {len(chunk['bullets'])} 要点, {len(chunk['numbered_lists'])} 数字列表")
            else:
                # 纯标题页
                title_layout = layout_map.get('TitleOnly', list(layout_map.values())[0])
                slide = prs.slides.add_slide(title_layout)
                add_content_to_slide_enhanced(slide, {'plain_text': [], 'bullets': [], 'numbered_lists': [], 'tables': [], 'code_blocks': [], 'quotes': []}, main_title)
                total_slides_created += 1
                print(f"  创建一级标题页: {main_title}")
        
        # 2. 处理各个section
        for section in sections:
            section_title = section['title']
            section_content = section['content']
            subsections = section.get('subsections', [])
            
            # 合并section内容和subsections内容
            all_section_content = section_content.copy()
            if MERGE_SUBSECTIONS and subsections:
                for subsection in subsections:
                    all_section_content.extend(subsection['content'])
            else:
                # 为每个subsection创建独立幻灯片
                for subsection in subsections:
                    sub_elements = extract_content_elements_enhanced(subsection['content'])
                    
                    # 检查是否存在混合内容（要点+代码块）
                    has_bullets_or_text = len(sub_elements['bullets']) > 0 or len(sub_elements['numbered_lists']) > 0 or len(sub_elements['plain_text']) > 0
                    has_code_blocks = len(sub_elements['code_blocks']) > 0
                    
                    if has_bullets_or_text and has_code_blocks:
                        print(f"    检测到子章节混合内容，将按原始顺序创建幻灯片")
                        
                        # 使用新的按顺序分割函数
                        ordered_chunks = split_content_by_order(sub_elements)
                        
                        for chunk_idx, chunk in enumerate(ordered_chunks):
                            layout = choose_layout_enhanced(chunk)
                            slide = prs.slides.add_slide(layout)
                            
                            chunk_title = get_chunk_title(subsection['title'], chunk, chunk_idx, len(ordered_chunks))
                            
                            add_content_to_slide_enhanced(slide, chunk, chunk_title, "")
                            total_slides_created += 1
                            
                            content_desc = []
                            if chunk['bullets']:
                                content_desc.append(f"{len(chunk['bullets'])} 要点")
                            if chunk['numbered_lists']:
                                content_desc.append(f"{len(chunk['numbered_lists'])} 数字列表")
                            if chunk['plain_text']:
                                content_desc.append(f"{len(chunk['plain_text'])} 文本段落")
                            if chunk['code_blocks']:
                                content_desc.append(f"{len(chunk['code_blocks'])} 代码块")
                            
                            print(f"      创建子章节混合内容页 {chunk_idx+1}/{len(ordered_chunks)}: {', '.join(content_desc)}")
                    else:
                        # 非混合内容，使用按顺序分割以保持原始顺序
                        sub_chunks = split_content_by_order(sub_elements)
                        
                        for chunk_idx, chunk in enumerate(sub_chunks):
                            layout = choose_layout_enhanced(chunk)
                            slide = prs.slides.add_slide(layout)
                            
                            subtitle = f"第{chunk_idx+1}部分" if len(sub_chunks) > 1 else ""
                            add_content_to_slide_enhanced(slide, chunk, subsection['title'], subtitle)
                            total_slides_created += 1
                            print(f"    创建子章节页 {chunk_idx+1}/{len(sub_chunks)}: {subsection['title']}")
            
            # 处理section主要内容
            if all_section_content:
                elements = extract_content_elements_enhanced(all_section_content)
                
                # 检查是否存在混合内容（要点+代码块）
                has_bullets_or_text = len(elements['bullets']) > 0 or len(elements['numbered_lists']) > 0 or len(elements['plain_text']) > 0
                has_code_blocks = len(elements['code_blocks']) > 0
                
                if has_bullets_or_text and has_code_blocks:
                    print(f"  检测到混合内容，将按原始顺序创建幻灯片")
                    
                    # 使用新的按顺序分割函数
                    ordered_chunks = split_content_by_order(elements)
                    
                    for chunk_idx, chunk in enumerate(ordered_chunks):
                        layout = choose_layout_enhanced(chunk)
                        slide = prs.slides.add_slide(layout)
                        
                        chunk_title = get_chunk_title(section_title, chunk, chunk_idx, len(ordered_chunks))
                        
                        add_content_to_slide_enhanced(slide, chunk, chunk_title, "")
                        total_slides_created += 1
                        
                        content_desc = []
                        if chunk['bullets']:
                            content_desc.append(f"{len(chunk['bullets'])} 要点")
                        if chunk['numbered_lists']:
                            content_desc.append(f"{len(chunk['numbered_lists'])} 数字列表")
                        if chunk['plain_text']:
                            content_desc.append(f"{len(chunk['plain_text'])} 文本段落")
                        if chunk['code_blocks']:
                            content_desc.append(f"{len(chunk['code_blocks'])} 代码块")
                        
                        print(f"    创建混合内容页 {chunk_idx+1}/{len(ordered_chunks)}: {', '.join(content_desc)}")
                else:
                    # 非混合内容，使用按顺序分割以保持原始顺序
                    chunks = split_content_by_order(elements)
                    
                    for chunk_idx, chunk in enumerate(chunks):
                        layout = choose_layout_enhanced(chunk)
                        slide = prs.slides.add_slide(layout)
                        
                        subtitle = f"第{chunk_idx+1}部分" if len(chunks) > 1 else ""
                        add_content_to_slide_enhanced(slide, chunk, section_title, subtitle)
                        total_slides_created += 1
                        print(f"  创建章节页 {chunk_idx+1}/{len(chunks)}: {len(chunk['plain_text'])} 文本段落, {len(chunk['bullets'])} 要点, {len(chunk['numbered_lists'])} 数字列表")
    
    return total_slides_created
