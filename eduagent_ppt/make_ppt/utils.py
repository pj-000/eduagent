"""
通用工具函数
提供各种辅助功能
"""

from pptx.util import Pt
from config import *


def apply_paragraph_format_respecting_template(paragraph, text_frame_original_paragraph=None):
    """
    应用段落格式，但尊重模板中的设置
    优先使用模板设置，只有在模板没有设置时才使用默认值
    """
    if text_frame_original_paragraph is not None:
        # 尝试获取并应用原始格式
        try:
            # 行间距
            if hasattr(text_frame_original_paragraph, 'line_spacing') and text_frame_original_paragraph.line_spacing is not None:
                paragraph.line_spacing = text_frame_original_paragraph.line_spacing
            
            # 段后间距
            if hasattr(text_frame_original_paragraph, 'space_after') and text_frame_original_paragraph.space_after is not None:
                paragraph.space_after = text_frame_original_paragraph.space_after
              # 段前间距 - 强制设置为0，不使用模板的段前间距
            paragraph.space_before = Pt(0)
                
            # 对齐方式
            if hasattr(text_frame_original_paragraph, 'alignment') and text_frame_original_paragraph.alignment is not None:
                paragraph.alignment = text_frame_original_paragraph.alignment
                
        except Exception as e:
            print(f"  警告: 获取模板段落格式时出错: {e}")
            # 如果获取失败，不设置任何格式，让PowerPoint使用默认值

def apply_default_content_format(paragraph, content_type="normal"):
    """
    在没有模板格式时应用默认格式
    """
    # 所有内容类型都设置段前间距为0
    paragraph.space_before = Pt(0)
    
    if content_type == "normal":
        # 普通文本使用适中的行距和间距
        paragraph.line_spacing = 1.3
        paragraph.space_after = Pt(8)
    elif content_type == "bullet":
        # 要点使用较紧凑的行距
        paragraph.line_spacing = 1.2  
        paragraph.space_after = Pt(3)
    elif content_type == "numbered":
        # 数字列表使用较紧凑的行距
        paragraph.line_spacing = 1.2
        paragraph.space_after = Pt(3)
    elif content_type == "quote":
        # 引用使用适中的行距
        paragraph.line_spacing = 1.2
        paragraph.space_after = Pt(6)