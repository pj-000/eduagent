"""
内容格式化工具
处理代码高亮、表格格式化、文本格式化等
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from config import *
from utils import apply_paragraph_format_respecting_template
import re


def apply_code_syntax_highlighting(paragraph, code_line):
    """
    为代码行应用语法高亮，并根据前导空格设置缩进
    改进版：更好地处理字符串和各种语法元素（不处理注释）
    """
    paragraph.text = ""  # 清空段落文本
    
    # 如果是空行，添加空行并返回
    if not code_line.strip():
        run = paragraph.add_run()
        run.text = ""  # 空行保持格式
        run.font.name = 'Consolas'
        run.font.size = Pt(14)
        return

    # 直接使用原始代码行（包含缩进）进行语法高亮
    i = 0
    
    while i < len(code_line):
        # 处理字符串（双引号和单引号）
        if code_line[i] in ['"', "'"]:
            quote = code_line[i]
            j = i + 1
            # 找到字符串结束，处理转义字符
            while j < len(code_line):
                if code_line[j] == quote:
                    # 检查是否为转义引号
                    backslash_count = 0
                    k = j - 1
                    while k >= i and code_line[k] == '\\':
                        backslash_count += 1
                        k -= 1
                    # 如果反斜杠数量为偶数，则引号未被转义
                    if backslash_count % 2 == 0:
                        j += 1  # 包含结束引号
                        break
                j += 1
            
            string_text = code_line[i:j]
            run = paragraph.add_run()
            run.text = string_text
            run.font.name = 'Consolas'
            run.font.size = Pt(14)
            run.font.color.rgb = STRING_COLOR
            i = j
            continue
        
        # 处理数字
        if code_line[i].isdigit():
            j = i
            while j < len(code_line) and (code_line[j].isdigit() or code_line[j] == '.'):
                j += 1
            
            number_text = code_line[i:j]
            run = paragraph.add_run()
            run.text = number_text
            run.font.name = 'Consolas'
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(0, 102, 204)  # 蓝色数字
            i = j
            continue
        
        # 处理标识符和关键字
        if code_line[i].isalpha() or code_line[i] == '_':
            j = i
            while j < len(code_line) and (code_line[j].isalnum() or code_line[j] == '_'):
                j += 1

            word = code_line[i:j]
            run = paragraph.add_run()
            run.text = word
            run.font.name = 'Consolas'
            run.font.size = Pt(14)

            # 应用关键字高亮
            if word in PYTHON_KEYWORDS:
                run.font.color.rgb = KEYWORD_COLOR
                run.font.bold = True
            else:
                # 普通标识符使用默认颜色
                run.font.color.rgb = RGBColor(0, 0, 0)  # 黑色

            i = j
            continue
        
        # 处理运算符和特殊字符
        if code_line[i] in '+-*/=<>!&|()[]{},.;:':
            run = paragraph.add_run()
            run.text = code_line[i]
            run.font.name = 'Consolas'
            run.font.size = Pt(14)
            run.font.color.rgb = RGBColor(128, 0, 128)  # 紫色运算符
            i += 1
            continue
        
        # 处理空格和制表符（保持原始格式）
        if code_line[i] in [' ', '\t']:
            space_start = i
            while i < len(code_line) and code_line[i] in [' ', '\t']:
                i += 1
            space_text = code_line[space_start:i]
            run = paragraph.add_run()
            run.text = space_text
            run.font.name = 'Consolas'
            run.font.size = Pt(14)
            continue
          # 处理其他字符（包括#符号，但不作为注释处理）
        run = paragraph.add_run()
        run.text = code_line[i]
        run.font.name = 'Consolas'
        run.font.size = Pt(14)
        run.font.color.rgb = RGBColor(0, 0, 0)  # 默认黑色
        i += 1

def distribute_code_to_placeholders(code_content, content_placeholders, layout_name):
    """
    将代码内容分布到多个内容占位符中
    """
    code_lines = code_content.split('\n')
    total_lines = len(code_lines)
    
    if layout_name == 'CodeSlide2' and len(content_placeholders) >= 2:
        # CodeSlide2 有两个占位符，将代码分为两部分
        mid_point = total_lines // 2
        
        # 第一个占位符
        text_frame1 = content_placeholders[0].text_frame
        original_paragraph1 = text_frame1.paragraphs[0] if text_frame1.paragraphs else None
        text_frame1.clear()
        
        for line_num, code_line in enumerate(code_lines[:mid_point]):
            if line_num == 0:
                p = text_frame1.paragraphs[0]
            else:
                p = text_frame1.add_paragraph()
            
            apply_code_syntax_highlighting(p, code_line)
            apply_paragraph_format_respecting_template(p, original_paragraph1)
        
        # 第二个占位符
        text_frame2 = content_placeholders[1].text_frame
        original_paragraph2 = text_frame2.paragraphs[0] if text_frame2.paragraphs else None
        text_frame2.clear()
        
        for line_num, code_line in enumerate(code_lines[mid_point:]):
            if line_num == 0:
                p = text_frame2.paragraphs[0]
            else:
                p = text_frame2.add_paragraph()
            
            apply_code_syntax_highlighting(p, code_line)
            apply_paragraph_format_respecting_template(p, original_paragraph2)
        
        print(f"  代码分布: 第一部分 {mid_point} 行，第二部分 {total_lines - mid_point} 行")
    
    else:
        # CodeSlide1 或只有一个占位符，使用第一个占位符
        text_frame = content_placeholders[0].text_frame
        original_paragraph = text_frame.paragraphs[0] if text_frame.paragraphs else None
        text_frame.clear()
        
        for line_num, code_line in enumerate(code_lines):
            if line_num == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            apply_code_syntax_highlighting(p, code_line)
            apply_paragraph_format_respecting_template(p, original_paragraph)
        
        print(f"  代码全部放入第一个占位符: {total_lines} 行")

def format_table_content(table_content, text_frame):
    """
    改进表格格式化，使用更好的行距和字体，并遵循占位符样式
    """
    # 保存原始格式信息
    original_line_spacing = None
    original_space_after = None
    original_space_before = None
    original_font_size = Pt(18)  # 默认字体大小
    
    # 尝试从占位符获取默认样式
    try:
        if text_frame.paragraphs and text_frame.paragraphs[0].runs:
            original_font_size = text_frame.paragraphs[0].runs[0].font.size or Pt(18)
        if text_frame.paragraphs:
            original_line_spacing = text_frame.paragraphs[0].line_spacing
            original_space_after = text_frame.paragraphs[0].space_after
            original_space_before = text_frame.paragraphs[0].space_before
    except:
        pass  # 使用默认值
    
    text_frame.clear()
    
    table_lines = table_content.split('\n')
    filtered_lines = []
    
    for line in table_lines:
        if '---' not in line and line.strip():
            # 处理表格行，改善格式
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if cells:
                formatted_line = ' | '.join(cells)
                filtered_lines.append(formatted_line)
    
    for i, line in enumerate(filtered_lines):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = line
        p.alignment = PP_ALIGN.LEFT
        
        # 设置表格样式 - 遵循占位符的字体大小
        if i == 0:  # 表头
            for run in p.runs:
                run.font.bold = True
                run.font.size = original_font_size
                run.font.name = 'Consolas'
        else:  # 表格内容
            for run in p.runs:
                run.font.size = Pt(max(16, int(original_font_size.pt * 0.9)))
                run.font.name = 'Consolas'
        
        # 使用模板的行间距设置，如果没有则使用默认值
        if original_line_spacing is not None:
            p.line_spacing = original_line_spacing
        if original_space_after is not None:
            p.space_after = original_space_after
        else:
            p.space_after = Pt(6)  # 只在模板没有设置时使用默认值
            
        # 强制设置段前间距为0，不使用模板的段前间距
        p.space_before = Pt(0)



def format_text_with_bold(paragraph, text):
    """
    处理包含粗体标记的文本，正确设置粗体格式
    """
    paragraph.text = ""  # 清空段落文本
    
    if '**' in text:
        # 按粗体标记分割文本
        parts = re.split(r'(\*\*[^*]+\*\*)', text)
        for part in parts:
            if part.startswith('**') and part.endswith('**') and len(part) > 4:
                # 粗体文本
                bold_text = part[2:-2]  # 移除**标记
                run = paragraph.add_run()
                run.text = bold_text
                run.font.bold = True
            elif part.strip():
                # 普通文本
                run = paragraph.add_run()
                run.text = part

    else:
        # 没有粗体标记的普通文本
        paragraph.text = text
        # 保持模板默认字号